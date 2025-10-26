#!/usr/bin/env python3
"""
A script to generate a short defense PPTX for the "Smart Home Voice Control" project.

This script uses python-pptx and requests to create a short presentation with placeholder images.
It produces presentation/short_defense.pptx in the repository when run.

Usage:
  pip install python-pptx requests Pillow
  python presentation/generate_short_ppt.py

You can later replace the placeholder images in the generated PPTX or edit the script to use local images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import requests
from io import BytesIO
import os

# Output path
OUTPUT_DIR = "presentation"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "short_defense.pptx")

# Placeholder images (via.placeholder.com)
PLACEHOLDER_IMAGE_URLS = {
    "architecture": "https://via.placeholder.com/1200x675.png?text=系统架构(占位图)",
    "hardware": "https://via.placeholder.com/1200x675.png?text=硬件示意(占位图)",
    "demo": "https://via.placeholder.com/1200x675.png?text=演示输出(占位图)",
}

SLIDES = [
    {
        "title": "基于语言识别的智能家居控制系统",
        "bullets": [
            "项目实训答辩 - 简短版",
            "学生：李耀均  学号：22216920122",
            "指导老师：李文莉",
            "广东东软学院 计算机学院"
        ],
        "image": None,
    },
    {
        "title": "目录",
        "bullets": ["项目背景与目标", "系统架构与关键功能", "核心模块(ASR/ NLU/ 设备)", "演示与运行方式", "问题与改进", "总结"],
        "image": None,
    },
    {
        "title": "项目背景与目标",
        "bullets": [
            "智能家居市场增长，语音交互为主流入口",
            "目标：实现语音→意图→设备控制的闭环",
            "支持灯光、空调、窗帘等基本控制与场景联动"
        ],
        "image": None,
    },
    {
        "title": "系统总体架构",
        "bullets": ["语音采集端：麦克风/手机", "控制中心：FastAPI + Whisper + NLU", "设备层：SmartLight、MQTT/HTTP/继电器/红外"],
        "image": "architecture",
    },
    {
        "title": "关键模块：ASR / NLU / 设备",
        "bullets": ["ASR：OpenAI Whisper（本地/云）", "NLU：规则+关键字/实体抽取", "设备控制：SmartDevice 抽象、SmartLight 实现"],
        "image": None,
    },
    {
        "title": "硬件与运行方式",
        "bullets": ["推荐：树莓派4/5 + USB麦克风", "继电器、红外发射、智能灯/插座", "启动：python start_server.py --host 0.0.0.0 --port 8000 --reload"],
        "image": "hardware",
    },
    {
        "title": "演示与运行（可现场演示）",
        "bullets": ["示例脚本：examples/demo_without_asr.py", "示例指令：打开客厅灯、将书房灯设为蓝色", "API 文档：/docs"],
        "image": "demo",
    },
    {
        "title": "目前成果与存在问题",
        "bullets": ["已实现：ASR、NLU、设备抽象、API 服务、示例脚本", "问题：噪声/方言识别、离线模型资源消耗大、协议兼容性需扩展"],
        "image": None,
    },
    {
        "title": "总结与致谢",
        "bullets": ["短期：优化NLU规则与测试", "中期：引入轻量模型与唤醒功能", "感谢指导老师与听众，欢迎提问"],
        "image": None,
    }
]


def download_image(url: str) -> BytesIO:
    """Download image from URL and return BytesIO object."""
    resp = requests.get(url, timeout=10)
    resp.raise_for_status()
    return BytesIO(resp.content)


def add_title_and_bullets(slide, title: str, bullets: list[str]):
    """Add title and bullet text to a slide."""
    title_shape = slide.shapes.title
    title_shape.text = title

    # add body
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(4.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = bullets[0] if bullets else ""
    p.font.size = Pt(18)
    p.level = 0

    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0


def create_presentation(output_path: str):
    prs = Presentation()

    # set slide width/height for 16:9
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for s in SLIDES:
        slide_layout = prs.slide_layouts[5]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        # add title and bullets
        add_title_and_bullets(slide, s["title"], s["bullets"]) 

        # if image specified, download and add
        if s.get("image"):
            key = s["image"]
            url = PLACEHOLDER_IMAGE_URLS.get(key)
            try:
                img_stream = download_image(url)
                left = Inches(7.5)
                top = Inches(1.2)
                width = Inches(5)
                slide.shapes.add_picture(img_stream, left, top, width=width)
            except Exception as e:
                print(f"警告：无法下载占位图 ({url}): {e}")

    # ensure output dir
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"生成演示文稿: {output_path}")


if __name__ == '__main__':
    create_presentation(OUTPUT_FILE)
